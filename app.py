import os
import re
import tempfile
from html import escape

import PyPDF2
import streamlit as st
from docx import Document
from pptx import Presentation
from pptx.util import Pt

st.set_page_config(page_title="CV to PPT Generator", layout="wide")


SECTION_ALIASES = {
    "profile": {
        "profile",
        "summary",
        "professional summary",
        "career summary",
        "objective",
        "about me",
    },
    "experience": {
        "experience",
        "work experience",
        "professional experience",
        "employment history",
        "relevant experience",
        "projects",
        "project",
    },
    "skills": {
        "skills",
        "technical skills",
        "core skills",
        "key skills",
        "competencies",
        "core competencies",
        "technical expertise",
        "functional technical experience",
        "functional and technical experience",
        "tools knowledge",
        "tool knowledge",
    },
    "education": {
        "education",
        "academic qualification",
        "academics",
        "qualification",
        "certification",
        "certifications",
    },
}

ACTION_WORDS = {
    "achieved",
    "analyzed",
    "automated",
    "built",
    "collaborated",
    "contributed",
    "created",
    "delivered",
    "designed",
    "developed",
    "drove",
    "enhanced",
    "executed",
    "implemented",
    "improved",
    "led",
    "managed",
    "optimized",
    "performed",
    "reduced",
    "responsible",
    "streamlined",
    "supported",
    "tested",
    "worked",
}

MONTH_PATTERN = (
    r"jan|feb|mar|apr|may|jun|jul|aug|sep|sept|oct|nov|dec|present|current"
)
CONTACT_PATTERN = re.compile(
    r"(@|linkedin|github|www\.|http|mailto:|\b\d{10,}\b)", re.IGNORECASE
)


def normalize_line(line):
    line = (
        line.replace("\xa0", " ")
        .replace("\uf0b7", "-")
        .replace("•", "-")
        .replace("\t", " ")
    )
    line = re.sub(r"\s+", " ", line).strip(" |-")
    return line


def split_rich_line(line):
    parts = re.split(r"\s*[|]\s*|\s{2,}", line)
    return [normalize_line(part) for part in parts if normalize_line(part)]


def is_probable_heading(line):
    cleaned = re.sub(r"[^a-z& ]", "", line.lower()).strip()
    if not cleaned:
        return False

    words = cleaned.split()
    if len(words) > 6:
        return False

    known_aliases = {alias for aliases in SECTION_ALIASES.values() for alias in aliases}
    if cleaned in known_aliases:
        return True

    return line.isupper()


def detect_section(line):
    cleaned = re.sub(r"[^a-z& ]", "", line.lower()).strip()
    for section, aliases in SECTION_ALIASES.items():
        if cleaned in aliases:
            return section
    return None


def is_contact_line(line):
    return bool(CONTACT_PATTERN.search(line))


def is_date_heavy(line):
    return bool(re.search(rf"\b({MONTH_PATTERN})\b|\b(19|20)\d{{2}}\b", line, re.IGNORECASE))


def extract_text_from_docx(file):
    doc = Document(file)
    chunks = [para.text for para in doc.paragraphs if para.text.strip()]

    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                chunks.append(" | ".join(cells))

    return "\n".join(chunks)


def extract_text_from_pdf(file):
    text = []
    reader = PyPDF2.PdfReader(file)
    for page in reader.pages:
        extracted = page.extract_text() or ""
        if extracted.strip():
            text.append(extracted)
    return "\n".join(text)


def extract_text(uploaded_file):
    if uploaded_file.name.endswith(".docx"):
        return extract_text_from_docx(uploaded_file)
    if uploaded_file.name.endswith(".pdf"):
        return extract_text_from_pdf(uploaded_file)
    return ""


def collect_lines(text):
    raw_lines = [normalize_line(line) for line in text.splitlines()]
    lines = []
    for raw_line in raw_lines:
        if not raw_line:
            continue
        parts = split_rich_line(raw_line)
        if parts:
            lines.extend(parts)
        else:
            lines.append(raw_line)
    return lines


def choose_name(lines):
    for line in lines[:12]:
        if is_contact_line(line) or is_probable_heading(line) or is_date_heavy(line):
            continue
        if len(line.split()) < 2 or len(line.split()) > 5:
            continue
        if re.search(r"[^A-Za-z.\- ]", line):
            continue
        return line
    return "Candidate Name"


def choose_location(lines):
    for line in lines[:15]:
        if is_contact_line(line):
            continue
        if re.search(r"\d", line):
            continue
        if "," in line and len(line.split()) <= 6:
            return line
    return ""


def choose_job_title(lines, name, location):
    ignored_lines = {name.strip().lower(), location.strip().lower()}

    for line in lines[:20]:
        lower_line = line.lower().strip()
        if not lower_line or lower_line in ignored_lines:
            continue
        if is_contact_line(line) or is_probable_heading(line) or is_date_heavy(line):
            continue
        if len(line) < 6 or len(line) > 60:
            continue
        return line

    return "Automation Testing or Functional Testing"


def split_sections(lines):
    sections = {key: [] for key in SECTION_ALIASES}
    current_section = None

    for line in lines:
        detected_section = detect_section(line)
        if detected_section:
            current_section = detected_section
            continue

        if is_probable_heading(line):
            current_section = None
            continue

        if current_section:
            sections[current_section].append(line)

    return sections


def clean_sentence(line, max_len=120):
    line = re.sub(r"\s+", " ", line).strip(" -")
    if len(line) <= max_len:
        return line

    trimmed = line[:max_len].rsplit(" ", 1)[0].strip()
    return f"{trimmed}..."


def score_experience_line(line):
    lower_line = line.lower()
    score = 0

    if any(word in lower_line for word in ACTION_WORDS):
        score += 4
    if re.search(r"\b\d+%|\b\d+[+,]?\d*\b", line):
        score += 2
    if "client" in lower_line or "project" in lower_line:
        score += 1
    if is_contact_line(line):
        score -= 10
    if is_probable_heading(line):
        score -= 10
    if len(line) < 20:
        score -= 2

    return score


def top_ranked_lines(lines, limit, scorer, max_len=120):
    candidates = []
    seen = set()

    for index, line in enumerate(lines):
        cleaned = clean_sentence(line, max_len=max_len)
        lowered = cleaned.lower()
        if lowered in seen:
            continue
        seen.add(lowered)
        candidates.append((scorer(cleaned), index, cleaned))

    selected = sorted(candidates, key=lambda item: (-item[0], item[1]))[:limit]
    selected.sort(key=lambda item: item[1])
    return [item[2] for item in selected if item[0] >= 0]


def build_profile_text(sections, lines):
    source_lines = sections["profile"] or lines[:20]
    profile_lines = []

    for line in source_lines:
        if is_contact_line(line) or is_probable_heading(line):
            continue
        if len(line) < 30:
            continue
        profile_lines.append(clean_sentence(line, max_len=105))

    if not profile_lines:
        profile_lines = top_ranked_lines(lines, limit=4, scorer=score_experience_line, max_len=105)

    return "\n".join(f"- {line}" for line in profile_lines[:4])


def build_experience_text(sections, lines):
    source_lines = sections["experience"] or lines
    experience_lines = top_ranked_lines(source_lines, limit=5, scorer=score_experience_line, max_len=110)
    return "\n".join(f"- {line}" for line in experience_lines)


def extract_skill_items(sections, lines):
    source_lines = sections["skills"] or lines
    items = []
    seen = set()

    for line in source_lines:
        for piece in re.split(r",|;|:|\(|\)", line):
            item = normalize_line(piece)
            if not item:
                continue
            if len(item) > 35 or len(item) < 2:
                continue
            if is_contact_line(item) or is_date_heavy(item):
                continue

            words = item.split()
            if len(words) > 5:
                continue

            lower_item = item.lower()
            if lower_item in seen:
                continue
            if lower_item in {"skills", "technical skills", "certifications"}:
                continue

            seen.add(lower_item)
            items.append(item)

    return items[:12]


def build_skills_text(sections, lines):
    skill_items = extract_skill_items(sections, lines)
    grouped_lines = []

    for start in range(0, len(skill_items), 3):
        chunk = skill_items[start : start + 3]
        if chunk:
            grouped_lines.append("- " + ", ".join(chunk))

    return "\n".join(grouped_lines[:4])


def build_education_text(sections):
    education_lines = top_ranked_lines(
        sections["education"], limit=3, scorer=score_experience_line, max_len=100
    )
    return "\n".join(f"- {line}" for line in education_lines)


def parse_cv(text):
    lines = collect_lines(text)
    sections = split_sections(lines)
    name = choose_name(lines)
    location = choose_location(lines)

    return {
        "name": name,
        "location": location,
        "job_title": choose_job_title(lines, name, location),
        "profile": build_profile_text(sections, lines),
        "experience": build_experience_text(sections, lines),
        "skills": build_skills_text(sections, lines),
        "education": build_education_text(sections),
    }


def text_shapes(slide):
    return [shape for shape in slide.shapes if getattr(shape, "has_text_frame", False)]


def shape_text(shape):
    return (shape.text or "").strip()


def find_shape_with_text(slide, marker):
    marker = marker.lower()
    for shape in text_shapes(slide):
        if marker in shape_text(shape).lower():
            return shape
    return None


def find_top_banner_shape(slide):
    candidates = []
    for shape in text_shapes(slide):
        if shape.top > 1200000:
            continue
        if shape.width < 6000000:
            continue
        candidates.append(shape)

    if not candidates:
        return None

    return min(candidates, key=lambda shape: (shape.top, shape.left))


def horizontal_overlap(shape_a, shape_b):
    left = max(shape_a.left, shape_b.left)
    right = min(shape_a.left + shape_a.width, shape_b.left + shape_b.width)
    return max(0, right - left)


def has_meaningful_horizontal_overlap(shape_a, shape_b, threshold=0.6):
    overlap = horizontal_overlap(shape_a, shape_b)
    baseline = min(shape_a.width, shape_b.width)
    if baseline <= 0:
        return False
    return (overlap / baseline) >= threshold


def find_body_shape_below_heading(slide, heading_text):
    heading_shape = find_shape_with_text(slide, heading_text)
    if not heading_shape:
        return None

    candidates = []
    for shape in text_shapes(slide):
        if shape == heading_shape:
            continue
        if shape.top <= heading_shape.top:
            continue
        if not has_meaningful_horizontal_overlap(shape, heading_shape):
            continue
        candidates.append(shape)

    if not candidates:
        return None

    return min(candidates, key=lambda shape: shape.top - heading_shape.top)


def find_name_shape(slide):
    named_shape = find_shape_with_text(slide, "FirstName SecondName")
    if named_shape:
        return named_shape

    profile_heading = find_shape_with_text(slide, "PROFILE")
    title_shape = find_top_banner_shape(slide)

    candidates = []
    for shape in text_shapes(slide):
        if shape == profile_heading or shape == title_shape:
            continue
        if title_shape and shape.top <= title_shape.top:
            continue
        if profile_heading and shape.top >= profile_heading.top:
            continue
        if shape.left > 4000000:
            continue
        candidates.append(shape)

    if not candidates:
        return None

    return max(candidates, key=lambda shape: shape.width * shape.height)


def set_paragraph_font(paragraph, size, bold=False):
    for run in paragraph.runs:
        run.font.size = Pt(size)
        run.font.bold = bold


def set_text_block(shape, content, font_size):
    if not shape:
        return

    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True

    lines = [line.strip() for line in content.splitlines() if line.strip()]
    if not lines:
        return

    first_paragraph = frame.paragraphs[0]
    first_paragraph.text = lines[0]
    first_paragraph.space_after = Pt(2)
    set_paragraph_font(first_paragraph, font_size)

    for line in lines[1:]:
        paragraph = frame.add_paragraph()
        paragraph.text = line
        paragraph.space_after = Pt(2)
        set_paragraph_font(paragraph, font_size)


def set_name_block(shape, name, location):
    if not shape:
        return

    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True

    name_paragraph = frame.paragraphs[0]
    name_paragraph.text = name or "Candidate Name"
    set_paragraph_font(name_paragraph, 22, bold=True)

    if location:
        location_paragraph = frame.add_paragraph()
        location_paragraph.text = location
        set_paragraph_font(location_paragraph, 11)


def set_single_line_text(shape, content, font_size, bold=False):
    if not shape:
        return

    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True

    text = (content or "").strip()
    if not text:
        return

    paragraph = frame.paragraphs[0]
    paragraph.text = text
    set_paragraph_font(paragraph, font_size, bold=bold)


def locate_template_shapes(slide):
    return {
        "job_title": find_shape_with_text(slide, "Automation Testing or Functional Testing")
        or find_top_banner_shape(slide),
        "name": find_name_shape(slide),
        "profile": find_body_shape_below_heading(slide, "PROFILE"),
        "experience": find_body_shape_below_heading(slide, "RELEVANT EXPERIENCE"),
        "skills": find_body_shape_below_heading(slide, "FUNCTIONAL & TECHNICAL EXPERIENCE"),
    }


def clear_template_content(shape_map):
    for shape in shape_map.values():
        if shape and getattr(shape, "has_text_frame", False):
            shape.text_frame.clear()
            shape.text_frame.word_wrap = True


def create_ppt(data, template_path):
    prs = Presentation(template_path)

    for slide in prs.slides:
        shape_map = locate_template_shapes(slide)
        clear_template_content(shape_map)
        set_single_line_text(shape_map["job_title"], data["job_title"], 20)
        set_name_block(shape_map["name"], data["name"], data["location"])
        set_text_block(shape_map["profile"], data["profile"], 11)
        set_text_block(shape_map["experience"], data["experience"], 10)
        set_text_block(shape_map["skills"], data["skills"], 10)

    output_path = os.path.join(tempfile.gettempdir(), "generated.pptx")
    prs.save(output_path)
    return output_path


def format_preview_lines(text):
    lines = [line.strip() for line in text.splitlines() if line.strip()]
    if not lines:
        return "<div style='color:#9ca3af;'>No content</div>"

    rendered = []
    for line in lines:
        cleaned = escape(line.lstrip("- ").strip())
        rendered.append(
            "<div style='margin-bottom:8px; line-height:1.45; color:#111827;'>"
            f"{cleaned}</div>"
        )
    return "".join(rendered)


def render_preview(data):
    preview_html = f"""
    <div style="background:linear-gradient(135deg, #eef4f8, #dde7ef); padding:18px; border-radius:18px;">
        <div style="background:#ffffff; width:100%; aspect-ratio:16 / 9; margin:0 auto; padding:18px 20px; box-shadow:0 16px 40px rgba(15, 23, 42, 0.12); border:1px solid #dbe4ea; overflow:hidden;">
            <div style="display:flex; justify-content:space-between; align-items:flex-start; margin-bottom:12px;">
                <div style="font-size:12px; color:#4b5563;">Template Preview</div>
                <div style="font-size:26px; font-weight:800; color:#111827; letter-spacing:-1px;">Deloitte<span style="color:#84b819;">.</span></div>
            </div>
            <div style="font-size:14px; color:#4b5563; margin-bottom:12px;">{escape(data["job_title"] or "Automation Testing or Functional Testing")}</div>
            <div style="display:grid; grid-template-columns:42% 58%; gap:10px; height:calc(100% - 78px);">
                <div style="display:grid; grid-template-rows:auto 1fr; gap:10px;">
                    <div style="border-top:3px solid #2d90c7; padding-top:10px;">
                        <div style="font-size:18px; font-weight:700; color:#111827; line-height:1.2;">{escape(data["name"] or "Candidate Name")}</div>
                        <div style="font-size:11px; color:#6b7280; margin-top:4px;">{escape(data["location"] or "Location")}</div>
                    </div>
                    <div>
                        <div style="background:#f97316; color:#ffffff; font-weight:700; padding:5px 8px; font-size:11px;">PROFILE</div>
                        <div style="border:1px solid #cfd8df; padding:10px; height:100%; font-size:11px; overflow:hidden;">
                            {format_preview_lines(data["profile"])}
                        </div>
                    </div>
                </div>
                <div style="display:grid; grid-template-rows:58% 42%; gap:10px;">
                    <div>
                        <div style="background:#f97316; color:#ffffff; font-weight:700; padding:5px 8px; font-size:11px;">RELEVANT EXPERIENCE</div>
                        <div style="border:1px solid #cfd8df; padding:10px; height:100%; font-size:10px; overflow:hidden;">
                            {format_preview_lines(data["experience"])}
                        </div>
                    </div>
                    <div>
                        <div style="background:#f97316; color:#ffffff; font-weight:700; padding:5px 8px; font-size:11px;">FUNCTIONAL &amp; TECHNICAL EXPERIENCE</div>
                        <div style="border:1px solid #cfd8df; padding:10px; height:100%; font-size:10px; overflow:hidden;">
                            {format_preview_lines(data["skills"])}
                        </div>
                    </div>
                </div>
            </div>
        </div>
    </div>
    """
    st.markdown(preview_html, unsafe_allow_html=True)


st.markdown(
    """
    <style>
    .main .block-container {
        max-width: 96rem;
        padding-top: 2rem;
        padding-left: 2rem;
        padding-right: 2rem;
    }

    div[data-testid="column"]:first-child {
        min-width: 0;
    }

    div[data-testid="column"]:last-child {
        min-width: 0;
    }
    </style>
    """,
    unsafe_allow_html=True,
)

st.title("CV to PPT Generator")

uploaded_file = st.file_uploader("Upload CV (PDF/DOCX)", type=["pdf", "docx"])

if uploaded_file:
    st.success("File uploaded successfully.")

    file_signature = f"{uploaded_file.name}-{uploaded_file.size}"
    if st.session_state.get("loaded_file_signature") != file_signature:
        text = extract_text(uploaded_file)
        parsed_data = parse_cv(text)
        st.session_state["loaded_file_signature"] = file_signature
        for field in ["name", "location", "job_title", "profile", "experience", "skills", "education"]:
            st.session_state[f"cv_{field}"] = parsed_data.get(field, "")

    editor_column, preview_column = st.columns([1, 2], gap="large")

    with editor_column:
        st.subheader("Edit Content")
        st.text_input("Job Title", key="cv_job_title")
        st.text_input("Name", key="cv_name")
        st.text_input("Location", key="cv_location")
        st.text_area("Profile", key="cv_profile", height=190)
        st.text_area("Relevant Experience", key="cv_experience", height=220)
        st.text_area("Functional / Technical Experience", key="cv_skills", height=190)

    edited_data = {
        "job_title": st.session_state.get("cv_job_title", "").strip(),
        "name": st.session_state.get("cv_name", "").strip(),
        "location": st.session_state.get("cv_location", "").strip(),
        "profile": st.session_state.get("cv_profile", "").strip(),
        "experience": st.session_state.get("cv_experience", "").strip(),
        "skills": st.session_state.get("cv_skills", "").strip(),
        "education": st.session_state.get("cv_education", "").strip(),
    }

    with preview_column:
        st.subheader("Preview")
        render_preview(edited_data)

    ppt_path = create_ppt(edited_data, "Deloitte Profile format 2026.pptx")

    with open(ppt_path, "rb") as file_handle:
        st.download_button(
            label="Download PPT",
            data=file_handle,
            file_name="Generated_Profile.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
